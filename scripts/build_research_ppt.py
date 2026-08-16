from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT=Path(__file__).resolve().parents[1]; OUT=ROOT/'AMR水印攻击实验汇报.pptx'
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
NAVY=RGBColor(18,35,58); BLUE=RGBColor(39,105,180); CYAN=RGBColor(45,170,190); ORANGE=RGBColor(235,145,45); RED=RGBColor(201,70,70); GREEN=RGBColor(47,145,95); LIGHT=RGBColor(242,246,250); MID=RGBColor(98,112,128); WHITE=RGBColor(255,255,255); DARK=RGBColor(30,38,48)

def bg(slide,color=WHITE):
 s=slide.background.fill; s.solid(); s.fore_color.rgb=color
def box(slide,x,y,w,h,text='',fill=LIGHT,line=None,fs=20,bold=False,color=DARK,align=PP_ALIGN.LEFT,radius=True):
 sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line or fill
 tf=sh.text_frame; tf.clear(); tf.margin_left=tf.margin_right=Inches(.18); tf.margin_top=tf.margin_bottom=Inches(.1); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
 p=tf.paragraphs[0]; p.text=text; p.alignment=align; p.font.name='Microsoft YaHei'; p.font.size=Pt(fs); p.font.bold=bold; p.font.color.rgb=color
 return sh
def title(slide,t,sub=None):
 box(slide,.55,.25,12.2,.62,t,WHITE,WHITE,27,True,NAVY,PP_ALIGN.LEFT,False)
 box(slide,.55,.87,1.0,.07,'',CYAN,CYAN,1,False,CYAN,PP_ALIGN.LEFT,False)
 if sub: box(slide,.55,1.0,12,.35,sub,WHITE,WHITE,12,False,MID,PP_ALIGN.LEFT,False)
def footer(slide,n): box(slide,12.35,7.12,.5,.22,str(n),WHITE,WHITE,10,False,MID,PP_ALIGN.RIGHT,False)
def bullets(slide,items,x=.75,y=1.45,w=11.8,h=5.5,fs=22):
 sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=sh.text_frame; tf.word_wrap=True
 for j,it in enumerate(items):
  p=tf.paragraphs[0] if j==0 else tf.add_paragraph(); p.text=it; p.level=0; p.font.name='Microsoft YaHei'; p.font.size=Pt(fs); p.font.color.rgb=DARK; p.space_after=Pt(13); p.text='• '+p.text
 return sh
def slide(): s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); return s

# 1
s=slide(); bg(s,NAVY); box(s,.75,1.2,11.8,1.2,'面向 AMR 语义水印检测的\n反馈式对抗攻击研究',NAVY,NAVY,34,True,WHITE); box(s,.78,3.0,10.8,.6,'研究目标 · 方法演进 · 正式实验 · 质量风险 · 下一步',NAVY,NAVY,18,False,RGBColor(190,215,235)); box(s,.78,5.65,4.8,.55,'RealNews · 250段 · 1000候选',CYAN,CYAN,18,True,WHITE,PP_ALIGN.CENTER); box(s,9.6,6.55,2.8,.35,'实验汇报｜2026.07',NAVY,NAVY,12,False,WHITE,PP_ALIGN.RIGHT); footer(s,1)
# 2
s=slide(); title(s,'1. 我们想做什么？','研究问题与核心假设')
box(s,.75,1.5,5.8,4.8,'目标\n\n在尽量保持原文语义、实体、数字、否定和可读性的前提下，改变文本的 AMR 结构，使 SWAN 水印检测分数下降。',LIGHT,BLUE,24,True,NAVY)
box(s,6.8,1.5,5.75,2.15,'核心问题',WHITE,WHITE,20,True,NAVY); box(s,7.1,2.15,5.05,1.25,'能否系统性规避\n基于 AMR 模板匹配的检测器？',RGBColor(231,242,252),BLUE,25,True,BLUE,PP_ALIGN.CENTER)
box(s,6.8,4.05,5.75,2.25,'攻击成功 = 检测规避 + 语义质量',WHITE,WHITE,20,True,NAVY); box(s,7.1,4.75,2.25,.85,'z ≤ 2.33',RGBColor(226,246,239),GREEN,22,True,GREEN,PP_ALIGN.CENTER); box(s,9.75,4.75,2.25,.85,'质量合格',RGBColor(255,241,224),ORANGE,22,True,ORANGE,PP_ALIGN.CENTER); footer(s,2)
# 3
s=slide(); title(s,'2. 我们发现了什么问题？','早期实验暴露出两个层面的脆弱性')
box(s,.75,1.55,5.8,4.85,'检测器层面',RGBColor(231,242,252),BLUE,24,True,NAVY); bullets(s,['SWAN依赖AMR概念、角色和子图模板匹配','句界变化、谓词重组会改变解析结果','单次简单改写通常只能让少量段落降分'],1.05,2.25,5.1,3.8,19)
box(s,6.8,1.55,5.75,4.85,'攻击质量层面',RGBColor(255,241,224),ORANGE,24,True,NAVY); bullets(s,['低z-score不等于语义保持','常见失败：否定反转、关系模糊、信息删减/新增','自动生成模型同时评分会产生自评偏差'],7.1,2.25,5.0,3.8,19); footer(s,3)
# 4
s=slide(); title(s,'3. 参考了什么？','方法与工具基础')
refs=[('SWAN','AMR语义水印与模板匹配检测'),('AMR / BART parser','把自然语言重新解析为语义图'),('S2MATCH / 图匹配','度量概念、角色与子图匹配'),('DeepSeek API','生成整段改写与攻击候选'),('RealNews','正式评估语料：250个完整段落')]
for i,(a,b) in enumerate(refs):
 y=1.45+i*1.02; box(s,.85,y,2.6,.72,a,NAVY,NAVY,18,True,WHITE,PP_ALIGN.CENTER); box(s,3.7,y,8.65,.72,b,LIGHT,LIGHT,19,False,DARK)
footer(s,4)
# 5
s=slide(); title(s,'4. 攻击路线如何演进？','从单算子试探到反馈式候选搜索')
steps=[('E1','跨语言往返','低分率 16.7–20.0%'),('E2','整段多候选 + SWAN反馈','形成当前主攻击'),('E3','跨语言 + AMR扰动组合','未明显超过E2'),('增强','质量感知 + 自适应补救','降低语义失败')]
for i,(a,b,c) in enumerate(steps):
 x=.65+i*3.15; box(s,x,1.8,2.65,.62,a,BLUE,BLUE,20,True,WHITE,PP_ALIGN.CENTER); box(s,x,2.55,2.65,1.3,b,LIGHT,BLUE,18,True,NAVY,PP_ALIGN.CENTER); box(s,x,4.05,2.65,.95,c,WHITE,WHITE,15,False,MID,PP_ALIGN.CENTER)
 if i<3: box(s,x+2.72,3.02,.38,.13,'',CYAN,CYAN,1)
box(s,1.15,5.65,11.0,.7,'关键发现：最强因素不是某一个句法变换，而是“生成多个候选 + 检测器反馈选优”',RGBColor(226,246,239),GREEN,20,True,GREEN,PP_ALIGN.CENTER); footer(s,5)
# 6
s=slide(); title(s,'5. 当前最强攻击怎么做？','E2：质量感知的检测器反馈搜索')
flow=[('① 原段落','250段'),('② 生成候选','每段4个'),('③ AMR重解析','1000个图组'),('④ 质量过滤','事实/实体/否定'),('⑤ SWAN评分','选择最低z')]
for i,(a,b) in enumerate(flow):
 x=.45+i*2.58; box(s,x,2.0,2.15,1.35,a+'\n'+b,LIGHT,BLUE,18,True,NAVY,PP_ALIGN.CENTER)
 if i<4: box(s,x+2.18,2.58,.34,.12,'',CYAN,CYAN,1)
box(s,1.1,4.35,11.1,1.25,'失败样本 → 定向追加4个候选 → 再解析 → 再过滤 → 再选优',RGBColor(255,241,224),ORANGE,22,True,ORANGE,PP_ALIGN.CENTER); footer(s,6)
# 7
s=slide(); title(s,'6. 实验设置','正式冻结配置')
cards=[('数据','RealNews\n250个完整段落'),('候选','4 × 250\n共1000个'),('解析','RTX 4060\nBART AMR parser'),('检测','SWAN\n阈值 z≤2.33'),('生成','DeepSeek API\n整段联合改写')]
for i,(a,b) in enumerate(cards):
 x=.55+i*2.55; box(s,x,1.75,2.15,2.0,a+'\n\n'+b,LIGHT,BLUE,18,True,NAVY,PP_ALIGN.CENTER)
bullets(s,['本地GPU负责AMR解析；CPU负责匹配与统计；API负责候选生成','显存不足时采用batch_size=4分块解析，不改变模型或实验定义','最终结果以段落级配对统计为单位'],.9,4.35,11.6,2.0,18); footer(s,7)
# 8
s=slide(); title(s,'7. 检测规避结果','攻击对SWAN区分能力造成显著下降')
metrics=[('99.6%','249/250\n选优候选低于阈值',GREEN),('0.728','AUROC\n基线约0.985',BLUE),('6.6%','TPR @ 1% FPR',ORANGE),('12.6%','TPR @ 5% FPR',RED)]
for i,(v,l,c) in enumerate(metrics):
 x=.7+i*3.12; box(s,x,1.65,2.65,1.25,v,c,c,30,True,WHITE,PP_ALIGN.CENTER); box(s,x,3.0,2.65,1.1,l,LIGHT,LIGHT,17,False,DARK,PP_ALIGN.CENTER)
img=ROOT/'experiments/planE_e2_joint250/detection/amr_roc_curve.png'
if img.exists(): s.shapes.add_picture(str(img),Inches(4.5),Inches(4.3),width=Inches(4.35),height=Inches(2.35))
box(s,.8,4.75,3.3,1.35,'结论\n检测规避能力出现质变',RGBColor(226,246,239),GREEN,22,True,GREEN,PP_ALIGN.CENTER); footer(s,8)
# 9
s=slide(); title(s,'8. 质量感知与自适应补救','为什么“先过滤质量，再选最低分”更合理')
vals=[('59.6%','先选最低z\n再检查质量'),('76.4%','先质量过滤\n再按z选优'),('87.2%','对失败段落\n自适应补救')]
for i,(v,l) in enumerate(vals):
 x=1.0+i*4.05; color=[MID,BLUE,GREEN][i]; box(s,x,1.75,3.25,1.25,v,color,color,32,True,WHITE,PP_ALIGN.CENTER); box(s,x,3.1,3.25,1.2,l,LIGHT,LIGHT,18,False,DARK,PP_ALIGN.CENTER)
box(s,1.15,5.05,11.0,1.0,'注意：87.2% 是宽松自动质量门下的结果，不是最终人工确认率',RGBColor(255,232,232),RED,20,True,RED,PP_ALIGN.CENTER); footer(s,9)
# 10
s=slide(); title(s,'9. 失败案例与质量风险','低检测分不等于有效攻击')
risks=[('否定反转','“没有宣布”被改成“确实介绍”'),('实体歧义','Washington 与 Washington, D.C.混合'),('新增信息','加入源文不存在的原因或归因'),('空泛退化','大量something / unspecified relation'),('可读性下降','长句堆叠、循环定义、模板化')]
for i,(a,b) in enumerate(risks):
 y=1.35+i*1.02; box(s,.75,y,2.25,.72,a,RED,RED,18,True,WHITE,PP_ALIGN.CENTER); box(s,3.25,y,9.1,.72,b,LIGHT,LIGHT,18,False,DARK)
footer(s,10)
# 11
s=slide(); title(s,'10. 人工复核告诉了我们什么？','结果对质量标准高度敏感')
box(s,.8,1.45,3.55,2.0,'宽松自动门\n\n218 / 250\n87.2%',RGBColor(226,246,239),GREEN,26,True,GREEN,PP_ALIGN.CENTER)
box(s,4.9,1.45,3.55,2.0,'保守模型二审\n\n17 / 250\n6.8%',RGBColor(255,241,224),ORANGE,26,True,ORANGE,PP_ALIGN.CENTER)
box(s,9.0,1.45,3.55,2.0,'分层60段首审\n\n7段严格有效\n5段同时规避',RGBColor(255,232,232),RED,23,True,RED,PP_ALIGN.CENTER)
box(s,1.0,4.25,11.3,1.45,'正确解读\n检测器确实被大幅削弱；但“保持语义的真实攻击成功率”尚不能只靠同一生成模型自动评分确定。',NAVY,NAVY,21,True,WHITE,PP_ALIGN.CENTER); footer(s,11)
# 12
s=slide(); title(s,'11. 我们目前得到的结论','可以说什么，暂时不能说什么')
box(s,.75,1.45,5.8,4.9,'可以明确说',RGBColor(226,246,239),GREEN,24,True,GREEN); bullets(s,['反馈式整段多候选搜索显著降低SWAN检测性能','250段正式实验中，选优后249段低于阈值','攻击效果不是由AMR解析失败造成','质量感知选择明显优于只追求最低z-score'],1.05,2.15,5.1,3.8,18)
box(s,6.8,1.45,5.75,4.9,'暂时不能直接说',RGBColor(255,232,232),RED,24,True,RED); bullets(s,['不能把99.6%写成语义保持攻击成功率','不能把87.2%当成人工确认结果','不能声称E3一定强于E2','不能仅凭单一RealNews数据证明跨领域普适性'],7.1,2.15,5.0,3.8,18); footer(s,12)
# 13
s=slide(); title(s,'12. 下一步怎么做？','把“强规避”推进为“可信论文结果”')
nexts=[('1','盲法人工评价','随机/分层抽取60–100段，至少两名评审'),('2','独立语义验证','NLI、语义相似度、实体与关系一致性'),('3','统计检验','置信区间、配对检验、攻击预算曲线'),('4','消融实验','候选数、反馈、AMR引导、自适应补救'),('5','跨域验证','增加另一个领域或数据集')]
for i,(n,a,b) in enumerate(nexts):
 y=1.25+i*1.05; box(s,.7,y,.6,.62,n,BLUE,BLUE,20,True,WHITE,PP_ALIGN.CENTER); box(s,1.5,y,2.5,.62,a,LIGHT,LIGHT,19,True,NAVY); box(s,4.15,y,8.1,.62,b,WHITE,WHITE,17,False,DARK)
footer(s,13)

prs.save(OUT); print(OUT, len(prs.slides))
